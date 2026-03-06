"""
generator.py — DeltaV Architecture PPT Drawing Engine
Generates ONE slide matching the DeltaV reference diagram style:
  - Room boxes with dashed borders
  - Rack cabinet frames (dark gray outer, light inner)
  - Stacked rack components in Server Cabinet
  - CHARM baseplate strips (orange/brown) in I/O Cabinets
  - Controller bus lines (blue)
  - Network bus line (magenta/pink)
  - Operator desk with monitors, workstation, printer
  - FOPP connectors between rooms
  - Software licenses strip
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Slide dimensions ──────────────────────────────────────────────────────────
SLIDE_W = 13.33
SLIDE_H = 7.50

# ── Color constants ───────────────────────────────────────────────────────────
def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C_CABINET_OUTER  = rgb('2B2B2B')   # dark rack frame
C_CABINET_INNER  = rgb('F0F0F0')   # light rack interior
C_CABINET_TOP    = rgb('3A3A3A')   # top/bottom cap of rack
C_RACK_COMP      = rgb('D8D8D8')   # generic rack component (light gray)
C_ROOM_BORDER    = rgb('888888')   # room border (simulated dashed)
C_ROOM_PDC_BG    = rgb('F9F9F9')   # PDC room fill
C_ROOM_OP_BG     = rgb('F9F9F9')   # Operator room fill
C_NETWORK_LINE   = rgb('CC0077')   # magenta network bus
C_CTRL_LINE      = rgb('0070C0')   # blue controller bus
C_CHARM_BASE     = rgb('8B4513')   # CHARM baseplate brown
C_CHARM_SLOT     = rgb('CC6600')   # CHARM card slot orange
C_CHARM_SLOT2    = rgb('E07800')   # lighter orange alternate
C_FOPP_FILL      = rgb('CC0077')   # FOPP connector magenta
C_DESK_TOP       = rgb('AAAAAA')   # desk surface
C_DESK_BODY      = rgb('CCCCCC')   # desk body
C_MONITOR_SCREEN = rgb('3A85C8')   # monitor screen (blue)
C_MONITOR_FRAME  = rgb('333333')   # monitor bezel
C_TOWER          = rgb('555555')   # tower workstation
C_PRINTER        = rgb('888888')   # printer
C_TEXT_DARK      = rgb('222222')
C_TEXT_MID       = rgb('444444')
C_WHITE          = rgb('FFFFFF')
C_UPS            = rgb('FFD700')
C_SWITCH         = rgb('2D6A2D')
C_WORKSTATION    = rgb('2B5EA7')
C_SCREEN_COMP    = rgb('555555')
C_SW_STRIP       = rgb('C8E6C9')   # software strip green
C_SW_BORDER      = rgb('4A8C3F')
C_FIREWALL       = rgb('C03030')


# ─────────────────────────────────────────────────────────────────────────────
# LOW-LEVEL SHAPE HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _rect(slide, x, y, w, h, fill=None, line_color=None, line_pt=0.75, no_fill=False):
    """Add a rectangle shape. All coords in inches."""
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if no_fill:
        s.fill.background()
    elif fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    s.text = ''
    return s


def _text(slide, text, x, y, w, h, size=9, bold=False, color=None,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box. All coords in inches."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = 'Calibri'
    if color:
        r.font.color.rgb = color
    return tb


def _line(slide, x1, y1, x2, y2, color, pt=1.5):
    """Draw a line between two points (inches)."""
    # Use a zero-height rectangle as a line (horizontal) or zero-width (vertical)
    if abs(y1 - y2) < 0.001:   # horizontal
        s = slide.shapes.add_shape(1, Inches(min(x1,x2)), Inches(y1 - 0.01),
                                   Inches(abs(x2-x1)), Inches(0.02))
    else:                        # vertical
        s = slide.shapes.add_shape(1, Inches(x1 - 0.01), Inches(min(y1,y2)),
                                   Inches(0.02), Inches(abs(y2-y1)))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.text = ''
    return s


def _set_dashed_border(shape, color, pt=1.5, dash='dash'):
    """Set a dashed line on a shape border via XML."""
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(pt)
    # Set dash type via XML
    spPr = shape._element.spPr
    lnEl = spPr.find(qn('a:ln'))
    if lnEl is None:
        return
    prstDash = etree.SubElement(lnEl, qn('a:prstDash'))
    prstDash.set('val', dash)   # 'dash', 'dashDot', 'sysDash'


# ─────────────────────────────────────────────────────────────────────────────
# RACK CABINET DRAWING
# ─────────────────────────────────────────────────────────────────────────────

CAP_H   = 0.14   # height of cabinet top/bottom cap
INNER_PAD = 0.06  # horizontal padding inside cabinet

# ── IO Cabinet tower sizing (R1 / R2 from architecture engine) ────────────────
CHARMS_PER_TOWER  = 8    # R1: baseplates per CIOC / tower
TOWERS_PER_IO_CAB = 3    # R2: max towers per IO cabinet
CHARMS_PER_IO_CAB = 24   # = 8 × 3


def _draw_rack_frame(slide, x, y, w, h):
    """Draw outer dark rack cabinet frame with top/bottom caps."""
    # Outer frame
    _rect(slide, x, y, w, h, fill=C_CABINET_OUTER, line_color=C_CABINET_OUTER, line_pt=1.5)
    # Inner light area
    _rect(slide, x + INNER_PAD, y + CAP_H,
          w - 2*INNER_PAD, h - 2*CAP_H,
          fill=C_CABINET_INNER, line_color=rgb('BBBBBB'), line_pt=0.5)
    # Top screw-strip detail
    _rect(slide, x + INNER_PAD, y + 0.03,
          w - 2*INNER_PAD, 0.06,
          fill=rgb('555555'), line_color=None)
    # Bottom screw-strip detail
    _rect(slide, x + INNER_PAD, y + h - 0.09,
          w - 2*INNER_PAD, 0.06,
          fill=rgb('555555'), line_color=None)


def _draw_server_cabinet(slide, x, y, w, h, items):
    """
    Draw a Server Cabinet with stacked rack-unit components.
    items: list of component dicts (sorted top to bottom)
    """
    _draw_rack_frame(slide, x, y, w, h)

    cx = x + INNER_PAD + 0.02
    cw = w - 2 * INNER_PAD - 0.04
    avail_h = h - 2 * CAP_H - 0.10
    start_y = y + CAP_H + 0.05

    # Define fixed-height items first, then distribute remainder
    ITEM_HEIGHTS = {
        'SCREEN':     0.30,
        'WORKSTATION':0.48,
        'UPS':        0.38,
        'SWITCH':     0.28,
        'FOPP':       0.26,
        'MEDIA_CONV': 0.26,
        'FIBER':      0.22,
    }

    # Calculate heights
    total_fixed = 0
    unknown_count = 0
    heights = []
    for item in items:
        dc = item.get('diagram_class', '')
        base = dc.replace('CHARM_', '').replace('_', '')
        h_item = None
        for k, v in ITEM_HEIGHTS.items():
            if k in dc:
                h_item = v
                break
        if h_item:
            total_fixed += h_item
        else:
            unknown_count += 1
        heights.append(h_item)

    remaining = max(avail_h - total_fixed, 0.10)
    default_h = (remaining / unknown_count) if unknown_count else 0.28
    default_h = max(0.20, min(0.45, default_h))

    cy = start_y
    for item, h_item in zip(items, heights):
        if h_item is None:
            h_item = default_h
        h_item = min(h_item, avail_h - (cy - start_y) - 0.02)
        if h_item < 0.12:
            break

        dc    = item.get('diagram_class', '')
        label = item.get('description', '')
        qty   = item.get('qty', 1)
        color = _comp_color(dc)

        # Component block
        comp = _rect(slide, cx, cy, cw, h_item - 0.03,
                     fill=color, line_color=rgb('999999'), line_pt=0.4)

        # Short label inside
        short = _short_label(label, qty)
        _text(slide, short, cx + 0.05, cy + 0.01, cw - 0.1, h_item - 0.06,
              size=6.5, bold=False, color=_text_color(color), wrap=False)

        cy += h_item


def _comp_color(diagram_class):
    cmap = {
        'SCREEN':     C_SCREEN_COMP, 'WORKSTATION': C_WORKSTATION,
        'UPS':        C_UPS,         'SWITCH':      C_SWITCH,
        'FOPP':       rgb('7B3F9E'), 'MEDIA_CONV':  rgb('00A0B0'),
        'FIBER':      rgb('6040A0'), 'CONTROLLER':  rgb('1A3A6B'),
        'CIOC':       rgb('2E5FA3'), 'POWER':       rgb('CC3300'),
        'FIREWALL':   C_FIREWALL,   'MONITOR':     rgb('4488AA'),
        'OPERATOR_WS':C_WORKSTATION,'CONSOLE':     rgb('5577AA'),
        'PRINTER':    rgb('888888'), 'SOFTWARE':    rgb('4A8C3F'),
    }
    for k, v in cmap.items():
        if k in diagram_class.upper():
            return v
    return C_RACK_COMP


def _text_color(bg_rgb):
    r, g, b = bg_rgb[0], bg_rgb[1], bg_rgb[2]
    brightness = (r * 299 + g * 587 + b * 114) / 1000
    return C_WHITE if brightness < 140 else C_TEXT_DARK


def _short_label(description: str, qty: int) -> str:
    """Create a concise label for inside a component block."""
    d = description.strip()
    # Shorten common patterns
    d = d.replace('DeltaV ', '').replace('Emerson ', '').replace('Assembly', 'Assy')
    d = d.replace('Redundant', 'Red.').replace('Workstation', 'WS')
    d = d.replace('Widescreen', '').replace(' Monitor', ' Mon.')
    if len(d) > 38:
        d = d[:36] + '…'
    prefix = f'[×{qty}]  ' if qty > 1 else ''
    return prefix + d


# ─────────────────────────────────────────────────────────────────────────────
# I/O CABINET DRAWING  (CHARM baseplates)
# ─────────────────────────────────────────────────────────────────────────────

def _draw_io_cabinet(slide, x, y, w, h, items):
    """
    Draw an I/O Cabinet with:
    - PK Controller(s) at top (PRI/SEC labels if redundant)
    - One CIOC label per tower
    - CHARM baseplate strips grouped into towers of 8
    - Tower separator lines between groups
    - Blue controller bus lines
    Uses exact baseplate count from items (R1 + R3).
    """
    _draw_rack_frame(slide, x, y, w, h)

    cx = x + INNER_PAD + 0.02
    cw = w - 2 * INNER_PAD - 0.04
    inner_y = y + CAP_H + 0.04
    inner_h = h - 2 * CAP_H - 0.06

    ctrl_items  = [it for it in items if it['diagram_class'] in ('CONTROLLER', 'CIOC')]
    charm_items = [it for it in items if it['diagram_class'] not in ('CONTROLLER', 'CIOC')]

    # ── Controller + CIOC blocks ───────────────────────────────────────────────
    ctrl_block_h = 0.0
    cy = inner_y
    for ctrl in ctrl_items:
        color = _comp_color(ctrl['diagram_class'])
        _rect(slide, cx, cy, cw, 0.30, fill=color,
              line_color=rgb('888888'), line_pt=0.5)
        _text(slide, _short_label(ctrl['description'], ctrl['qty']),
              cx + 0.05, cy + 0.02, cw - 0.1, 0.26,
              size=6.5, bold=True, color=C_WHITE, wrap=False)
        cy += 0.33
        ctrl_block_h += 0.33

    # ── Blue controller bus lines ─────────────────────────────────────────────
    if ctrl_items:
        bus_y     = cy + 0.02
        bus_left  = cx + cw * 0.2
        bus_right = cx + cw * 0.8
        _rect(slide, bus_left, bus_y, bus_right - bus_left, 0.015,
              fill=C_CTRL_LINE, line_color=None)
        for bx in [bus_left + 0.05, bus_right - 0.05]:
            _rect(slide, bx - 0.007, bus_y, 0.014, 0.06,
                  fill=C_CTRL_LINE, line_color=None)

    # ── R1: Exact CHARM baseplate count ───────────────────────────────────────
    baseplate_count = sum(
        it.get('qty', 1) for it in charm_items
        if 'CHARM' in it.get('diagram_class', '')
    )
    baseplate_count = max(4, baseplate_count)   # always show at least 4 visual slots

    # R1: One tower = 8 baseplates + 1 CIOC
    num_towers = max(1, math.ceil(baseplate_count / CHARMS_PER_TOWER))
    # R3: Cap at what one IO cabinet can hold
    num_strips = min(baseplate_count, CHARMS_PER_IO_CAB)

    charm_y = cy + (0.10 if ctrl_items else 0.04)
    charm_h = inner_h - ctrl_block_h - (0.14 if ctrl_items else 0.06)

    # ── Per-tower CIOC headers + baseplate strips ─────────────────────────────
    TOWER_SEP   = 0.05    # gap between towers
    strip_gap   = 0.016

    # Calculate strip width so all towers fit in cw
    total_sep_w = TOWER_SEP * (num_towers - 1)
    strip_w     = (cw - total_sep_w - strip_gap * (num_strips - 1)) / max(num_strips, 1)
    strip_w     = max(0.08, min(strip_w, 0.28))

    CIOC_H    = 0.17   # height of per-tower CIOC label bar
    CIOC_GAP  = 0.025

    sx = cx
    strips_drawn = 0

    for t in range(num_towers):
        strips_in_tower = min(CHARMS_PER_TOWER, num_strips - strips_drawn)
        if strips_in_tower <= 0:
            break

        tower_w = strips_in_tower * (strip_w + strip_gap) - strip_gap
        tower_x = sx

        # CIOC label bar for this tower
        _rect(slide, tower_x, charm_y, tower_w, CIOC_H,
              fill=rgb('2E5FA3'), line_color=rgb('1A3A6B'), line_pt=0.4)
        _text(slide, f'CIOC-{t + 1:02d}',
              tower_x, charm_y, tower_w, CIOC_H,
              size=5.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Baseplate strips for this tower
        bp_y = charm_y + CIOC_H + CIOC_GAP
        bp_h = charm_h - CIOC_H - CIOC_GAP

        for i in range(strips_in_tower):
            bx = tower_x + i * (strip_w + strip_gap)
            _rect(slide, bx, bp_y, strip_w, bp_h,
                  fill=C_CHARM_BASE, line_color=rgb('5A2D0C'), line_pt=0.4)
            slot_count = 16
            slot_h = (bp_h - 0.04) / slot_count
            for j in range(slot_count):
                slot_color = C_CHARM_SLOT if j % 2 == 0 else C_CHARM_SLOT2
                _rect(slide, bx + 0.012, bp_y + 0.02 + j * slot_h,
                      strip_w - 0.024, slot_h - 0.008,
                      fill=slot_color, line_color=None)

        # Tower separator (thin vertical rule between towers)
        if t < num_towers - 1:
            sep_x = tower_x + tower_w + TOWER_SEP / 2 - 0.005
            _rect(slide, sep_x, charm_y, 0.010, charm_h,
                  fill=rgb('CCCCCC'), line_color=None)

        sx += tower_w + TOWER_SEP
        strips_drawn += strips_in_tower

    # ── Label below CHARM area ────────────────────────────────────────────────
    tower_label = f'{num_towers} tower{"s" if num_towers > 1 else ""}'
    _text(slide, f'CHARM I/O  ·  {tower_label}  ·  {baseplate_count} baseplates',
          cx, charm_y + charm_h + 0.01, cw, 0.16,
          size=6.0, italic=True, color=C_TEXT_MID, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# OPERATOR DESK DRAWING
# ─────────────────────────────────────────────────────────────────────────────

def _draw_operator_desk(slide, x, y, w, h, items):
    """
    Draw operator desk with monitors, workstation tower, printer.
    Matches the reference: desk top surface, monitors on desk, tower beside.
    """
    has_printer = any('PRINTER' in it['diagram_class'] for it in items)
    has_firewall = any('FIREWALL' in it['diagram_class'] for it in items)

    # Reserve space for printer on right
    printer_w  = 0.60 if has_printer else 0
    desk_w     = w - printer_w - (0.08 if has_printer else 0)
    desk_x     = x

    # ── Firewall (small red bar at top-left) ──────────────────────────────────
    if has_firewall:
        _rect(slide, desk_x, y + 0.04, min(1.0, desk_w * 0.5), 0.22,
              fill=C_FIREWALL, line_color=rgb('880000'), line_pt=0.5)
        _text(slide, 'Firewall', desk_x + 0.03, y + 0.07, 0.90, 0.16,
              size=6.5, bold=True, color=C_WHITE)

    # ── OWS Label ─────────────────────────────────────────────────────────────
    _text(slide, 'OWS', desk_x + desk_w * 0.25, y + 0.04, desk_w * 0.5, 0.22,
          size=11, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)

    # ── Monitors on desk ──────────────────────────────────────────────────────
    mon_items  = [it for it in items if 'MONITOR' in it['diagram_class']]
    mon_count  = mon_items[0]['qty'] if mon_items else 2
    mon_count  = min(mon_count, 2)

    mon_w_total = desk_w * 0.80
    mon_w  = (mon_w_total - 0.06 * (mon_count - 1)) / mon_count
    mon_h  = mon_w * 0.60
    mon_y  = y + 0.30
    mon_start_x = desk_x + (desk_w - mon_w_total) / 2

    for i in range(mon_count):
        mx = mon_start_x + i * (mon_w + 0.06)
        _rect(slide, mx, mon_y, mon_w, mon_h,
              fill=C_MONITOR_FRAME, line_color=rgb('111111'), line_pt=0.5)
        _rect(slide, mx + 0.022, mon_y + 0.018,
              mon_w - 0.044, mon_h - 0.04,
              fill=C_MONITOR_SCREEN, line_color=None)
        # Stand
        _rect(slide, mx + mon_w * 0.38, mon_y + mon_h,
              mon_w * 0.24, 0.055, fill=C_MONITOR_FRAME, line_color=None)
        _rect(slide, mx + mon_w * 0.28, mon_y + mon_h + 0.05,
              mon_w * 0.44, 0.022, fill=C_MONITOR_FRAME, line_color=None)

    # ── Desk surface ───────────────────────────────────────────────────────────
    desk_surface_y = mon_y + mon_h + 0.08
    desk_surface_h = 0.10
    _rect(slide, desk_x, desk_surface_y, desk_w, desk_surface_h,
          fill=C_DESK_TOP, line_color=rgb('888888'), line_pt=0.5)

    # Desk body
    desk_body_h = max(0.60, h - (desk_surface_y - y) - desk_surface_h - 0.15)
    _rect(slide, desk_x, desk_surface_y + desk_surface_h,
          desk_w, desk_body_h,
          fill=C_DESK_BODY, line_color=rgb('999999'), line_pt=0.5)

    # ── Tower Workstation (inside desk body area on left side) ─────────────────
    twr_w  = min(desk_w * 0.40, 0.55)
    twr_h  = desk_body_h * 0.80
    twr_x  = desk_x + desk_w * 0.15
    twr_y  = desk_surface_y + desk_surface_h + (desk_body_h - twr_h) / 2
    _rect(slide, twr_x, twr_y, twr_w, twr_h,
          fill=C_TOWER, line_color=rgb('333333'), line_pt=0.5)
    _rect(slide, twr_x + 0.04, twr_y + twr_h * 0.08,
          twr_w - 0.08, twr_h * 0.08,
          fill=rgb('444444'), line_color=None)
    _text(slide, 'WS', twr_x + 0.02, twr_y + twr_h * 0.35,
          twr_w - 0.04, 0.20,
          size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── Printer (separate column on right) ────────────────────────────────────
    if has_printer:
        pr_x = x + desk_w + 0.08
        pr_y = y + 0.30
        pr_h = 0.42
        # Check bounds
        if pr_x + printer_w <= x + w + 0.02:
            _rect(slide, pr_x, pr_y, printer_w - 0.05, pr_h,
                  fill=C_PRINTER, line_color=rgb('555555'), line_pt=0.5)
            _rect(slide, pr_x + 0.06, pr_y + 0.10,
                  printer_w - 0.17, 0.04, fill=rgb('AAAAAA'), line_color=None)
            _text(slide, 'Printer', pr_x, pr_y + pr_h + 0.03,
                  printer_w - 0.05, 0.16,
                  size=6.5, color=C_TEXT_MID, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# FOPP CONNECTOR
# ─────────────────────────────────────────────────────────────────────────────

def _draw_fopp(slide, cx, cy, size=0.22):
    """Draw FOPP as a magenta diamond/arrow connector."""
    # Diamond shape (use a rotated square approximation with small shapes)
    half = size / 2
    # Left triangle
    tri_pts = [
        (cx - half, cy),
        (cx, cy - half),
        (cx, cy + half),
    ]
    # Draw as a filled rectangle rotated — approximate with two triangles
    # Use a simple rectangle with text "FOPP" as fallback for python-pptx
    s = slide.shapes.add_shape(1,
        Inches(cx - half), Inches(cy - half * 0.65),
        Inches(size), Inches(size * 0.65))
    s.fill.solid()
    s.fill.fore_color.rgb = C_FOPP_FILL
    s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'FOPP'
    r.font.size = Pt(5.5)
    r.font.bold = True
    r.font.color.rgb = C_WHITE


# ─────────────────────────────────────────────────────────────────────────────
# SOFTWARE STRIP
# ─────────────────────────────────────────────────────────────────────────────

def _draw_software_strip(slide, x, y, w, h, items):
    """Draw software licenses as a horizontal strip across PDC room top."""
    _rect(slide, x, y, w, h, fill=C_SW_STRIP,
          line_color=C_SW_BORDER, line_pt=0.8)

    labels = []
    for it in items:
        d = it['description']
        # Shorten
        for old, new in [('DeltaV ', ''), ('Emerson ', ''), ('Premium ', ''),
                         ('Workstation ', 'WS '), ('Software', 'SW'),
                         ('ProfessionalPLUS', 'ProPlus')]:
            d = d.replace(old, new)
        if len(d) > 28:
            d = d[:26] + '…'
        labels.append(d)

    text = '   |   '.join(labels)
    _text(slide, text, x + 0.12, y + 0.04, w - 0.24, h - 0.06,
          size=6.5, color=rgb('1B5E20'), wrap=False)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN SLIDE GENERATOR
# ─────────────────────────────────────────────────────────────────────────────

def generate_architecture_slide(structure: dict, project_title: str = 'DeltaV System Architecture') -> Presentation:
    """
    Generate a single-slide PPT architecture diagram.

    Args:
        structure: grouped BOM dict from grouper.group_bom()
        project_title: shown as the slide title

    Returns:
        python-pptx Presentation object (1 slide)
    """
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_WHITE

    # ── TITLE ─────────────────────────────────────────────────────────────────
    _text(slide, project_title,
          0.20, 0.06, SLIDE_W - 0.4, 0.30,
          size=16, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)
    # underline via a thin line
    _rect(slide, SLIDE_W * 0.3, 0.37, SLIDE_W * 0.4, 0.02, fill=C_TEXT_DARK)

    # ── LAYOUT CALCULATION ────────────────────────────────────────────────────
    pdc_room     = structure.get('PDC ROOM', {})
    op_room      = structure.get('OPERATOR ROOM', {})

    # Room boxes
    MARGIN_TOP   = 0.45
    MARGIN_BOT   = 0.18
    ROOM_Y       = MARGIN_TOP
    ROOM_H       = SLIDE_H - MARGIN_TOP - MARGIN_BOT

    PDC_X = 0.08
    PDC_W = 8.70
    OP_X  = PDC_X + PDC_W + 0.12
    OP_W  = SLIDE_W - OP_X - 0.08

    # ── PDC ROOM BOX ──────────────────────────────────────────────────────────
    pdc_box = _rect(slide, PDC_X, ROOM_Y, PDC_W, ROOM_H,
                    fill=C_ROOM_PDC_BG, line_color=C_ROOM_BORDER, line_pt=1.2)
    _set_dashed_border(pdc_box, C_ROOM_BORDER, pt=1.2, dash='sysDash')
    _text(slide, 'PDC ROOM', PDC_X + 0.15, ROOM_Y + 0.06, 3.0, 0.24,
          size=12, bold=False, color=C_TEXT_DARK)

    # ── OPERATOR ROOM BOX ────────────────────────────────────────────────────
    op_box = _rect(slide, OP_X, ROOM_Y, OP_W, ROOM_H,
                   fill=C_ROOM_OP_BG, line_color=C_ROOM_BORDER, line_pt=1.2)
    _set_dashed_border(op_box, C_ROOM_BORDER, pt=1.2, dash='sysDash')
    _text(slide, 'OPERATOR ROOM', OP_X + 0.15, ROOM_Y + 0.06, OP_W - 0.2, 0.24,
          size=12, bold=False, color=C_TEXT_DARK)

    # ── SOFTWARE STRIP ────────────────────────────────────────────────────────
    SW_H      = 0.26
    SW_Y      = ROOM_Y + 0.34
    sw_items  = pdc_room.get('SOFTWARE', [])
    if sw_items:
        _draw_software_strip(slide, PDC_X + 0.15, SW_Y, PDC_W - 0.30, SW_H, sw_items)

    # ── PDC CABINETS ──────────────────────────────────────────────────────────
    # Identify cabinet groups in PDC room (excluding SOFTWARE and FOPP_NODES)
    # FOPP_NODES are standalone topology connectors, drawn separately below
    pdc_cabinets = {k: v for k, v in pdc_room.items()
                    if k not in ('SOFTWARE', 'FOPP_NODES') and v}

    # Sort: SERVER_CABINET first, then I/O cabinets
    ordered_cabs = sorted(pdc_cabinets.keys(),
                          key=lambda k: (0 if 'SERVER' in k else 1, k))

    cab_top     = SW_Y + SW_H + 0.14 if sw_items else ROOM_Y + 0.38
    cab_h       = ROOM_H - (cab_top - ROOM_Y) - 0.65   # leave room for labels + FOPP
    cab_h       = max(3.5, cab_h)

    num_cabs    = len(ordered_cabs)
    total_cab_w = PDC_W - 0.35
    cab_gap     = 0.16

    # Server cabinet gets slightly less width than I/O cabinets
    server_w_ratio = 0.27
    io_w_ratio     = (1.0 - server_w_ratio) / max(num_cabs - 1, 1)

    widths = []
    for cab in ordered_cabs:
        if 'SERVER' in cab:
            widths.append(server_w_ratio)
        else:
            widths.append(io_w_ratio)
    # Normalize
    total = sum(widths)
    widths = [w / total for w in widths]

    usable_w  = total_cab_w - cab_gap * (num_cabs - 1)
    cab_x     = PDC_X + 0.17
    net_line_ys = []   # y-coord at top of each cabinet for network line

    for i, (cab_name, items) in enumerate(
            [(k, pdc_cabinets[k]) for k in ordered_cabs]):

        cw = usable_w * widths[i]
        cx = cab_x

        if 'SERVER' in cab_name:
            _draw_server_cabinet(slide, cx, cab_top, cw, cab_h, items)
        else:
            _draw_io_cabinet(slide, cx, cab_top, cw, cab_h, items)

        net_line_ys.append((cx + cw/2, cab_top + 0.02))   # top-center of cabinet

        # Cabinet label below
        _text(slide, cab_name.replace('_', ' ').title(),
              cx, cab_top + cab_h + 0.05, cw, 0.22,
              size=9, bold=False, color=C_TEXT_MID, align=PP_ALIGN.CENTER)

        cab_x += cw + cab_gap

    # ── NETWORK BUS LINE (magenta, horizontal across all cabinets) ───────────
    if net_line_ys:
        net_y    = cab_top - 0.12
        net_x1   = PDC_X + 0.17 + 0.08
        net_x2   = cab_x - cab_gap - 0.08
        # Horizontal pink line
        _rect(slide, net_x1, net_y, net_x2 - net_x1, 0.025,
              fill=C_NETWORK_LINE, line_color=None)
        # Vertical drops from each cabinet top to the network line
        for (cx_mid, _) in net_line_ys:
            _rect(slide, cx_mid - 0.010, net_y, 0.020, cab_top - net_y,
                  fill=C_NETWORK_LINE, line_color=None)

    # ── FOPP connectors + cross-room line ─────────────────────────────────────
    # R5 (other team): FOPP count comes from FOPP_NODES in structure, not hardcoded
    fopp_y = cab_top + cab_h + 0.28   # below cabinet labels

    fopp_items     = pdc_room.get('FOPP_NODES', [])
    pdc_fopp_count = max(1, len(fopp_items))   # always draw at least 1 FOPP

    # Space PDC FOPPs evenly across the right half of the PDC room
    fopp_spacing  = (PDC_W * 0.55) / (pdc_fopp_count + 1)
    pdc_fopp_xs   = [PDC_X + PDC_W * 0.45 + fopp_spacing * (i + 1)
                     for i in range(pdc_fopp_count)]

    for i, fx in enumerate(pdc_fopp_xs):
        _draw_fopp(slide, fx, fopp_y)
        fopp_lbl = fopp_items[i]['description'] if i < len(fopp_items) else 'FOPP'
        _text(slide, fopp_lbl[:10], fx - 0.25, fopp_y - 0.20, 0.50, 0.18,
              size=6.5, bold=True, color=C_TEXT_MID, align=PP_ALIGN.CENTER)
        # Pink line from network bus down to each FOPP
        _rect(slide, fx - 0.010, net_y, 0.020, fopp_y - net_y + 0.05,
              fill=C_NETWORK_LINE, line_color=None)

    # Use the rightmost PDC FOPP for the cross-room fiber connection
    pdc_fopp_x = pdc_fopp_xs[-1]

    # Operator side FOPP
    op_fopp_x = OP_X + 0.35
    _draw_fopp(slide, op_fopp_x, fopp_y)
    _text(slide, 'FOPP', op_fopp_x - 0.20, fopp_y - 0.20, 0.40, 0.18,
          size=6.5, bold=True, color=C_TEXT_MID, align=PP_ALIGN.CENTER)

    # Horizontal line between FOPPs
    fx1 = pdc_fopp_x + 0.12
    fx2 = op_fopp_x - 0.12
    _rect(slide, fx1, fopp_y - 0.010, fx2 - fx1, 0.020,
          fill=C_NETWORK_LINE, line_color=None)

    # Line from operator FOPP up to desk top area
    op_desk_connect_y = ROOM_Y + 0.60
    _rect(slide, op_fopp_x - 0.010, op_desk_connect_y, 0.020,
          fopp_y - op_desk_connect_y,
          fill=C_NETWORK_LINE, line_color=None)

    # ── OPERATOR DESK ────────────────────────────────────────────────────────
    desk_items = op_room.get('OPERATOR_DESK', [])
    desk_x  = OP_X + 0.12
    desk_y  = ROOM_Y + 0.34
    desk_w  = OP_W - 0.24
    desk_h  = cab_top + cab_h - desk_y + 0.24   # align bottom with cabinets
    desk_h  = max(3.8, min(desk_h, ROOM_H - 0.50))
    _draw_operator_desk(slide, desk_x, desk_y, desk_w, desk_h, desk_items)

    # ── LEGEND (bottom strip) ─────────────────────────────────────────────────
    legend_items = [
        ('Controller', rgb('1A3A6B')), ('CIOC',     rgb('2E5FA3')),
        ('CHARM I/O',  C_CHARM_SLOT),  ('Network',  C_SWITCH),
        ('EWS/WS',     C_WORKSTATION), ('UPS',      C_UPS),
        ('Fiber/FOPP', rgb('7B3F9E')), ('Firewall', C_FIREWALL),
        ('Power',      rgb('CC3300')),
    ]
    lx = PDC_X + 0.10
    ly = SLIDE_H - MARGIN_BOT + 0.02
    for lbl, col in legend_items:
        _rect(slide, lx, ly, 0.14, 0.12, fill=col,
              line_color=rgb('888888'), line_pt=0.3)
        _text(slide, lbl, lx + 0.17, ly - 0.01, 0.85, 0.14,
              size=5.5, color=C_TEXT_MID)
        lx += 1.04

    return prs


def generate_pptx(structure: dict, output_path: str,
                  project_title: str = 'DeltaV System Architecture'):
    """
    Full pipeline: structure → draw → save PPTX.
    """
    prs = generate_architecture_slide(structure, project_title)
    prs.save(output_path)
    return output_path