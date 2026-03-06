"""
app.py — DeltaV Architecture Generator
Aesthetic: Dark industrial-luxury control room
"""

import os, sys, io, tempfile, subprocess, shutil
import pandas as pd
import streamlit as st

sys.path.insert(0, os.path.dirname(__file__))
from parser     import parse_bom
from classifier import classify_dataframe, load_rules, save_user_correction
from grouper    import group_bom
from generator  import generate_pptx

st.set_page_config(
    page_title='DeltaV Architecture Generator',
    page_icon='⚡',
    layout='wide',
    initial_sidebar_state='expanded'
)

CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=JetBrains+Mono:wght@300;400;500&family=Inter:wght@300;400;500&display=swap');

*, *::before, *::after { box-sizing: border-box; }

html, body, [data-testid="stAppViewContainer"] {
    background: #0B0E14 !important;
    color: #E8EAF0 !important;
    font-family: 'Inter', sans-serif !important;
}

[data-testid="stAppViewContainer"] {
    background:
        radial-gradient(ellipse 80% 60% at 50% -10%, rgba(255,160,30,0.07) 0%, transparent 70%),
        repeating-linear-gradient(0deg, transparent, transparent 39px, rgba(255,255,255,0.018) 39px, rgba(255,255,255,0.018) 40px),
        repeating-linear-gradient(90deg, transparent, transparent 39px, rgba(255,255,255,0.018) 39px, rgba(255,255,255,0.018) 40px),
        #0B0E14 !important;
    background-attachment: fixed !important;
}

[data-testid="stSidebar"] {
    background: #0D1018 !important;
    border-right: 1px solid rgba(255,160,30,0.15) !important;
}
[data-testid="stSidebar"] * { color: #B8BCC8 !important; }

#MainMenu, footer, header { visibility: hidden; }
[data-testid="stDecoration"] { display: none; }
.block-container { padding: 0 2rem 2rem 2rem !important; max-width: 1400px !important; }

h1, h2, h3 { font-family: 'Syne', sans-serif !important; }

.hero { padding: 3rem 0 2.5rem 0; }
.hero-eyebrow {
    font-family: 'JetBrains Mono', monospace;
    font-size: 11px; font-weight: 500;
    letter-spacing: 0.22em; color: #FFA01E;
    text-transform: uppercase; margin-bottom: 0.8rem;
}
.hero-title {
    font-family: 'Syne', sans-serif;
    font-size: clamp(2rem, 5vw, 3.6rem);
    font-weight: 800; line-height: 1.05;
    color: #F0F2F8; margin: 0 0 1rem 0;
    letter-spacing: -0.02em;
}
.hero-title .accent { color: #FFA01E; }
.hero-sub {
    font-size: 15px; color: #7A8099;
    font-weight: 300; max-width: 580px; line-height: 1.7;
}
.hero-divider {
    height: 1px;
    background: linear-gradient(90deg, #FFA01E 0%, rgba(255,160,30,0.2) 40%, transparent 100%);
    margin: 2.5rem 0 2rem 0;
}

.metrics-row {
    display: grid;
    grid-template-columns: repeat(4, 1fr);
    gap: 1rem; margin-bottom: 2rem;
}
.metric-card {
    background: rgba(255,255,255,0.033);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 8px; padding: 1.2rem 1.4rem;
    position: relative; overflow: hidden;
}
.metric-card::before {
    content: ''; position: absolute;
    top: 0; left: 0; right: 0; height: 2px;
    background: linear-gradient(90deg, #FFA01E, transparent);
}
.metric-label {
    font-family: 'JetBrains Mono', monospace;
    font-size: 10px; letter-spacing: 0.15em;
    color: #5A6080; text-transform: uppercase; margin-bottom: 0.5rem;
}
.metric-value {
    font-family: 'Syne', sans-serif;
    font-size: 2.2rem; font-weight: 700;
    color: #F0F2F8; line-height: 1;
}
.metric-value.amber { color: #FFA01E; }
.metric-value.green { color: #4ADE80; }
.metric-value.red   { color: #F87171; }
.metric-sub { font-size: 11px; color: #5A6080; margin-top: 0.3rem; font-family: 'JetBrains Mono', monospace; }

.steps-row { display: flex; align-items: center; gap: 0; margin-bottom: 2rem; }
.step-item {
    display: flex; align-items: center; gap: 0.5rem;
    padding: 0.5rem 1.2rem;
    font-family: 'JetBrains Mono', monospace;
    font-size: 11px; letter-spacing: 0.08em;
    color: #404560; text-transform: uppercase;
}
.step-item.active { color: #FFA01E; }
.step-item.done   { color: #4ADE80; }
.step-num {
    width: 22px; height: 22px; border-radius: 50%;
    border: 1px solid currentColor;
    display: flex; align-items: center; justify-content: center;
    font-size: 10px; flex-shrink: 0;
}
.step-arrow { color: #252840; font-size: 18px; padding: 0 0.2rem; }

.section-header { display: flex; align-items: center; gap: 1rem; margin-bottom: 1rem; }
.section-num {
    font-family: 'JetBrains Mono', monospace; font-size: 10px;
    color: #FFA01E; letter-spacing: 0.1em;
    background: rgba(255,160,30,0.08);
    border: 1px solid rgba(255,160,30,0.2);
    border-radius: 4px; padding: 3px 8px; flex-shrink: 0;
}
.section-title { font-family: 'Syne', sans-serif; font-size: 18px; font-weight: 700; color: #E8EAF0; }
.section-line { flex: 1; height: 1px; background: rgba(255,255,255,0.06); }

[data-testid="stFileUploader"] > div {
    border: 1.5px dashed rgba(255,160,30,0.25) !important;
    background: rgba(255,160,30,0.025) !important;
    border-radius: 12px !important; padding: 1.5rem !important;
}
[data-testid="stFileUploader"] > div:hover {
    border-color: rgba(255,160,30,0.55) !important;
    background: rgba(255,160,30,0.05) !important;
}
[data-testid="stFileUploader"] label { color: #7A8099 !important; font-family: 'JetBrains Mono', monospace !important; font-size: 12px !important; }
[data-testid="stFileUploader"] button {
    background: rgba(255,160,30,0.12) !important;
    border: 1px solid rgba(255,160,30,0.35) !important;
    color: #FFA01E !important; border-radius: 6px !important;
    font-family: 'JetBrains Mono', monospace !important; font-size: 11px !important;
}

[data-testid="stDataFrame"] { background: transparent !important; }
[data-testid="stDataFrame"] th {
    background: rgba(255,160,30,0.08) !important; color: #FFA01E !important;
    font-family: 'JetBrains Mono', monospace !important; font-size: 10px !important;
    letter-spacing: 0.1em !important; text-transform: uppercase !important;
    border-bottom: 1px solid rgba(255,160,30,0.2) !important; padding: 10px 14px !important;
}
[data-testid="stDataFrame"] td {
    background: rgba(255,255,255,0.02) !important; color: #B8BCC8 !important;
    border-bottom: 1px solid rgba(255,255,255,0.04) !important;
    font-family: 'JetBrains Mono', monospace !important; font-size: 11px !important;
    padding: 8px 14px !important;
}

.cab-grid { display: grid; grid-template-columns: repeat(auto-fill, minmax(240px,1fr)); gap: 0.8rem; margin-top: 1rem; }
.cab-card {
    background: rgba(255,255,255,0.03);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 8px; padding: 1rem 1.2rem;
    position: relative; overflow: hidden;
}
.cab-card::before { content: ''; position: absolute; top:0;left:0;right:0;height:2px; }
.cab-card.pdc::before  { background: linear-gradient(90deg,#3B82F6,transparent); }
.cab-card.op::before   { background: linear-gradient(90deg,#4ADE80,transparent); }
.cab-card.sw::before   { background: linear-gradient(90deg,#A78BFA,transparent); }
.cab-room { font-family:'JetBrains Mono',monospace; font-size:9px; letter-spacing:0.15em; text-transform:uppercase; color:#5A6080; margin-bottom:0.3rem; }
.cab-name { font-family:'Syne',sans-serif; font-size:13px; font-weight:600; color:#D0D4E8; margin-bottom:0.5rem; }
.cab-count { font-family:'JetBrains Mono',monospace; font-size:22px; font-weight:700; color:#FFA01E; }
.cab-label { font-size:10px; color:#5A6080; margin-top:2px; }

[data-testid="stSelectbox"] > div {
    background: rgba(255,255,255,0.04) !important;
    border: 1px solid rgba(255,255,255,0.1) !important;
    border-radius: 6px !important; color: #E8EAF0 !important;
}
[data-testid="stTextInput"] input {
    background: rgba(255,255,255,0.04) !important;
    border: 1px solid rgba(255,255,255,0.1) !important;
    border-radius: 6px !important; color: #E8EAF0 !important;
    font-family: 'JetBrains Mono', monospace !important; font-size: 13px !important;
}
[data-testid="stTextInput"] input:focus { border-color: rgba(255,160,30,0.4) !important; }
[data-testid="stTextInput"] label {
    font-family: 'JetBrains Mono', monospace !important; font-size: 10px !important;
    letter-spacing: 0.1em !important; text-transform: uppercase !important; color: #5A6080 !important;
}

[data-testid="stButton"] button {
    background: linear-gradient(135deg,#FFA01E 0%,#FF7A00 100%) !important;
    color: #0B0E14 !important; border: none !important; border-radius: 8px !important;
    font-family: 'Syne', sans-serif !important; font-size: 14px !important;
    font-weight: 700 !important; letter-spacing: 0.04em !important;
    box-shadow: 0 4px 24px rgba(255,160,30,0.25) !important;
    text-transform: uppercase !important; transition: all 0.2s !important;
}
[data-testid="stButton"] button:hover { transform: translateY(-1px) !important; box-shadow: 0 8px 32px rgba(255,160,30,0.4) !important; }

[data-testid="stDownloadButton"] button {
    background: rgba(74,222,128,0.1) !important; color: #4ADE80 !important;
    border: 1px solid rgba(74,222,128,0.3) !important; border-radius: 8px !important;
    font-family: 'Syne', sans-serif !important; font-size: 14px !important;
    font-weight: 700 !important; width: 100% !important;
    text-transform: uppercase !important; transition: all 0.2s !important;
}
[data-testid="stDownloadButton"] button:hover { background: rgba(74,222,128,0.18) !important; box-shadow: 0 4px 20px rgba(74,222,128,0.2) !important; }

[data-testid="stExpander"] {
    background: rgba(255,255,255,0.02) !important;
    border: 1px solid rgba(255,255,255,0.07) !important; border-radius: 8px !important;
}

.output-card {
    background: linear-gradient(135deg,rgba(74,222,128,0.05) 0%,rgba(255,160,30,0.04) 100%);
    border: 1px solid rgba(74,222,128,0.2); border-radius: 12px;
    padding: 2rem; text-align: center; position: relative; overflow: hidden; margin-bottom: 1rem;
}
.output-card::before { content:''; position:absolute; top:0;left:0;right:0;height:2px; background:linear-gradient(90deg,#4ADE80,#FFA01E); }
.output-title { font-family:'Syne',sans-serif; font-size:20px; font-weight:700; color:#E8EAF0; margin-bottom:0.4rem; }
.output-sub { font-family:'JetBrains Mono',monospace; font-size:11px; color:#5A6080; margin-bottom:1.5rem; }

.sidebar-brand { font-family:'Syne',sans-serif; font-size:20px; font-weight:800; letter-spacing:-0.02em; }
.sidebar-ver { font-family:'JetBrains Mono',monospace; font-size:10px; letter-spacing:0.1em; color:#404560; margin-bottom:1.5rem; }
.sidebar-sec { font-family:'JetBrains Mono',monospace; font-size:9px; letter-spacing:0.18em; text-transform:uppercase; color:#404560; margin:1.2rem 0 0.6rem 0; padding-bottom:0.4rem; border-bottom:1px solid rgba(255,255,255,0.05); }
.tag-pill { display:inline-block; font-family:'JetBrains Mono',monospace; font-size:10px; padding:2px 8px; border-radius:20px; border:1px solid rgba(255,160,30,0.3); color:#FFA01E; background:rgba(255,160,30,0.07); margin:2px; }

::-webkit-scrollbar { width:6px; height:6px; }
::-webkit-scrollbar-track { background:transparent; }
::-webkit-scrollbar-thumb { background:rgba(255,160,30,0.2); border-radius:3px; }

/* ── DIAGRAM PREVIEW ── */
.preview-wrap {
    position: relative;
    border-radius: 12px;
    overflow: hidden;
    border: 1px solid rgba(255,160,30,0.18);
    box-shadow:
        0 0 0 1px rgba(255,255,255,0.04),
        0 24px 64px rgba(0,0,0,0.6),
        0 0 80px rgba(255,160,30,0.06);
    background: #141720;
    margin-bottom: 1.5rem;
}
.preview-topbar {
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 0.65rem 1rem;
    background: rgba(255,255,255,0.03);
    border-bottom: 1px solid rgba(255,255,255,0.06);
}
.preview-dots { display:flex; gap:6px; }
.preview-dot {
    width: 10px; height: 10px; border-radius: 50%;
}
.preview-dot-r { background: #FF5F57; }
.preview-dot-y { background: #FFBD2E; }
.preview-dot-g { background: #28CA41; }
.preview-label {
    font-family: 'JetBrains Mono', monospace;
    font-size: 10px;
    color: #404560;
    letter-spacing: 0.1em;
}
.preview-badge {
    font-family: 'JetBrains Mono', monospace;
    font-size: 9px;
    color: #4ADE80;
    background: rgba(74,222,128,0.08);
    border: 1px solid rgba(74,222,128,0.2);
    border-radius: 4px;
    padding: 2px 7px;
    letter-spacing: 0.08em;
}
.preview-img-wrap {
    padding: 0;
    background: #0D0F16;
}
.preview-footer {
    padding: 0.6rem 1rem;
    background: rgba(255,255,255,0.02);
    border-top: 1px solid rgba(255,255,255,0.04);
    display: flex;
    justify-content: space-between;
    align-items: center;
}
.preview-footer-left {
    font-family: 'JetBrains Mono', monospace;
    font-size: 10px;
    color: #2A2E40;
}
.preview-zoom-hint {
    font-family: 'JetBrains Mono', monospace;
    font-size: 10px;
    color: #2A2E40;
}
</style>
"""
st.markdown(CSS, unsafe_allow_html=True)


# ── SIDEBAR ──────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown('<div class="sidebar-brand" style="color:#FFA01E">⚡ DeltaV</div>', unsafe_allow_html=True)
    st.markdown('<div style="color:#C0C4D8;font-family:\'Syne\',sans-serif;font-size:13px;font-weight:600;margin-bottom:2px;">Architecture Generator</div>', unsafe_allow_html=True)
    st.markdown('<div class="sidebar-ver">v2.0 · EMERSON PUNE · FREE</div>', unsafe_allow_html=True)

    st.markdown('<div class="sidebar-sec">Project Settings</div>', unsafe_allow_html=True)
    project_title = st.text_input('PROJECT TITLE', value='CO2 Capture Plant')

    st.markdown('<div class="sidebar-sec">Supported Formats</div>', unsafe_allow_html=True)
    st.markdown('<span class="tag-pill">.xlsx</span><span class="tag-pill">.xls</span><span class="tag-pill">.csv</span>', unsafe_allow_html=True)

    st.markdown('<div class="sidebar-sec">Classification Levels</div>', unsafe_allow_html=True)
    st.markdown('<div style="font-family:\'JetBrains Mono\',monospace;font-size:11px;color:#5A6080;line-height:1.9;">① Part number prefix<br>② Exact phrase match<br>③ Keyword scoring<br>④ Category fallback</div>', unsafe_allow_html=True)

    st.markdown('<div class="sidebar-sec">Engine</div>', unsafe_allow_html=True)
    st.markdown('<div style="font-family:\'Inter\',sans-serif;font-size:12px;color:#5A6080;line-height:1.7;">100% rule-based · no API · fully offline · self-learning</div>', unsafe_allow_html=True)


# ── HERO ──────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="hero">
    <div class="hero-eyebrow">⬡ Emerson Automation Solutions · Pune · DCS Engineering</div>
    <h1 class="hero-title">BOM <span class="accent">→</span> Architecture<br>in Seconds</h1>
    <p class="hero-sub">Upload any Bill of Materials. The rule engine classifies every component, groups it into cabinets, and renders a one-slide DeltaV architecture diagram.</p>
</div>
<div class="hero-divider"></div>
""", unsafe_allow_html=True)


# ── STEPS ─────────────────────────────────────────────────────────────────────
def step_html(n, label, state):
    return f'<div class="step-item {state}"><div class="step-num">{n}</div>{label}</div>'

classified = 'classified_df' in st.session_state
generated  = 'pptx_bytes'   in st.session_state
s1 = 'done' if classified else 'active'
s2 = ('done' if generated else 'active') if classified else 'idle'
s3 = 'active' if generated else 'idle'

st.markdown(f"""
<div class="steps-row">
    {step_html(1,'Upload BOM',s1)}
    <span class="step-arrow">›</span>
    {step_html(2,'Classify & Review',s2)}
    <span class="step-arrow">›</span>
    {step_html(3,'Generate & Download',s3)}
</div>
""", unsafe_allow_html=True)


# ── STEP 1 ────────────────────────────────────────────────────────────────────
st.markdown('<div class="section-header"><span class="section-num">01</span><span class="section-title">Upload Bill of Materials</span><div class="section-line"></div></div>', unsafe_allow_html=True)

uploaded_file = st.file_uploader('Drop your BOM file — Excel or CSV, any format', type=['xlsx','xls','csv'])

if uploaded_file is None:
    st.markdown("""
    <div style="margin-top:1rem;padding:1.5rem;background:rgba(255,255,255,0.02);
         border:1px solid rgba(255,255,255,0.05);border-radius:8px;">
        <div style="font-family:'JetBrains Mono',monospace;font-size:11px;color:#404560;line-height:2.2;">
        ◦ &nbsp;Any column order — parser auto-detects Description, Qty, Area, Part No<br>
        ◦ &nbsp;Headers on any row (1–15), merged cells, section labels all handled<br>
        ◦ &nbsp;UTF-8, Latin-1, Windows-1252 encodings supported<br>
        ◦ &nbsp;Multi-sheet Excel — best sheet auto-selected
        </div>
    </div>
    """, unsafe_allow_html=True)
    st.stop()


# ── PARSE + CLASSIFY ──────────────────────────────────────────────────────────
suffix = os.path.splitext(uploaded_file.name)[1]
with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
    tmp.write(uploaded_file.read())
    tmp_path = tmp.name

with st.spinner('Parsing and classifying…'):
    try:
        df = parse_bom(tmp_path)
    except Exception as e:
        st.error(f'**Parse error:** {e}')
        st.stop()
    rules = load_rules()
    df    = classify_dataframe(df, rules)

unknowns = df[df['diagram_class'] == 'UNKNOWN']
known    = df[df['diagram_class'] != 'UNKNOWN']
pct      = int(len(known) / len(df) * 100) if len(df) else 0

# ── METRICS ───────────────────────────────────────────────────────────────────
st.markdown(f"""
<div class="metrics-row">
    <div class="metric-card">
        <div class="metric-label">Total Items</div>
        <div class="metric-value">{len(df)}</div>
        <div class="metric-sub">{uploaded_file.name[:22]}</div>
    </div>
    <div class="metric-card">
        <div class="metric-label">Classified</div>
        <div class="metric-value green">{len(known)}</div>
        <div class="metric-sub">{pct}% auto-resolved</div>
    </div>
    <div class="metric-card">
        <div class="metric-label">Unknowns</div>
        <div class="metric-value {'red' if len(unknowns) else 'green'}">{len(unknowns)}</div>
        <div class="metric-sub">{'need review' if len(unknowns) else 'all clear ✓'}</div>
    </div>
    <div class="metric-card">
        <div class="metric-label">PDC Items</div>
        <div class="metric-value amber">{len(df[df['area']=='PDC ROOM'])}</div>
        <div class="metric-sub">{len(df[df['area']=='OPERATOR ROOM'])} operator room</div>
    </div>
</div>
""", unsafe_allow_html=True)


# ── STEP 2 ────────────────────────────────────────────────────────────────────
st.markdown('<div class="section-header"><span class="section-num">02</span><span class="section-title">Review Classification</span><div class="section-line"></div></div>', unsafe_allow_html=True)

display_df = df[['sr_no','area','description','qty','diagram_class','confidence','part_number']].copy()
display_df.columns = ['SR','AREA','DESCRIPTION','QTY','CLASS','CONFIDENCE','PART NO']
st.dataframe(display_df, use_container_width=True, height=300)

corrections = {}
if len(unknowns) > 0:
    st.markdown(f'<div style="font-family:\'Syne\',sans-serif;font-size:15px;font-weight:600;color:#F87171;margin:1.2rem 0 0.8rem;">⚠ {len(unknowns)} item{"s" if len(unknowns)>1 else ""} need manual classification</div>', unsafe_allow_html=True)
    all_classes = sorted(rules.keys())
    for idx, row in unknowns.iterrows():
        c1, c2 = st.columns([3, 1])
        with c1:
            st.markdown(f'<div style="padding:0.5rem 0;font-size:13px;color:#E8EAF0;font-family:Inter,sans-serif;">{row["description"][:80]}<br><span style="font-family:JetBrains Mono,monospace;font-size:10px;color:#404560;">{row.get("part_number","")}</span></div>', unsafe_allow_html=True)
        with c2:
            choice = st.selectbox('', ['— select —'] + all_classes, key=f'fix_{idx}', label_visibility='collapsed')
            if choice != '— select —':
                corrections[idx] = choice

    if corrections:
        if st.button('💾  Save to rules.json', key='save_rules'):
            for idx, cls in corrections.items():
                save_user_correction(df.loc[idx, 'description'], cls)
            st.success(f'✅ Saved {len(corrections)} corrections — auto-classifies next time')
        for idx, cls in corrections.items():
            df.at[idx, 'diagram_class'] = cls
    else:
        df.loc[df['diagram_class'] == 'UNKNOWN', 'diagram_class'] = 'WORKSTATION'
else:
    st.success('✅  All items classified automatically — no review needed')

with st.expander('◈  Cabinet grouping preview'):
    structure = group_bom(df)
    room_cls  = {'PDC ROOM':'pdc', 'OPERATOR ROOM':'op'}
    html = '<div class="cab-grid">'
    for room, cabs in structure.items():
        for cab, items in cabs.items():
            if not items: continue
            rc = room_cls.get(room, 'sw')
            html += f'<div class="cab-card {rc}"><div class="cab-room">{room}</div><div class="cab-name">{cab.replace("_"," ").title()}</div><div class="cab-count">{len(items)}</div><div class="cab-label">components</div></div>'
    html += '</div>'
    st.markdown(html, unsafe_allow_html=True)


# ── STEP 3 ────────────────────────────────────────────────────────────────────
st.markdown('<div class="section-header" style="margin-top:2rem;"><span class="section-num">03</span><span class="section-title">Generate Architecture Diagram</span><div class="section-line"></div></div>', unsafe_allow_html=True)

col_btn, _ = st.columns([1, 2])
with col_btn:
    go = st.button('⚡  Generate Architecture PPT', type='primary', use_container_width=True)


def pptx_to_preview_image(pptx_bytes: bytes) -> bytes | None:
    """
    Render first slide of PPTX to PNG bytes using pure Python + Pillow.
    No LibreOffice or external tools required.
    Handles: filled rectangles, text boxes, connectors/lines.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image, ImageDraw, ImageFont
        import io as _io

        SCALE    = 96          # pixels per inch  (96 dpi → good balance of speed vs quality)
        CANVAS_W = int(13.33 * SCALE)
        CANVAS_H = int(7.50  * SCALE)

        prs   = Presentation(_io.BytesIO(pptx_bytes))
        slide = prs.slides[0]

        img  = Image.new('RGB', (CANVAS_W, CANVAS_H), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        def emu_to_px(emu):
            return int(emu / 914400 * SCALE)

        def get_rgb(color_obj):
            """Safely extract (r,g,b) from a pptx color object."""
            try:
                c = color_obj.rgb
                return (c[0], c[1], c[2])
            except Exception:
                return None

        def shape_fill(shape):
            try:
                f = shape.fill
                if f.type is None:
                    return None
                from pptx.enum.dml import MSO_THEME_COLOR
                rgb = f.fore_color.rgb
                return (rgb[0], rgb[1], rgb[2])
            except Exception:
                return None

        def shape_line_color(shape):
            try:
                ln = shape.line
                rgb = ln.color.rgb
                return (rgb[0], rgb[1], rgb[2])
            except Exception:
                return None

        def shape_line_width(shape):
            try:
                w = shape.line.width
                return max(1, int(w / 12700))   # EMU → rough px
            except Exception:
                return 1

        # ── Draw background ───────────────────────────────────────────────────
        try:
            bg = slide.background.fill
            bg_rgb = bg.fore_color.rgb
            img.paste(Image.new('RGB', (CANVAS_W, CANVAS_H),
                                (bg_rgb[0], bg_rgb[1], bg_rgb[2])))
        except Exception:
            pass

        # ── Draw shapes in z-order ────────────────────────────────────────────
        for shape in slide.shapes:
            try:
                x  = emu_to_px(shape.left   or 0)
                y  = emu_to_px(shape.top    or 0)
                w  = emu_to_px(shape.width  or 0)
                h  = emu_to_px(shape.height or 0)

                if w <= 0 or h <= 0:
                    continue

                fill   = shape_fill(shape)
                lcolor = shape_line_color(shape)
                lwidth = shape_line_width(shape)

                # ── Connector / line shape ────────────────────────────────────
                if shape.shape_type == 10:   # MSO_SHAPE_TYPE.LINE / CONNECTOR
                    try:
                        x1 = emu_to_px(shape.begin_x)
                        y1 = emu_to_px(shape.begin_y)
                        x2 = emu_to_px(shape.end_x)
                        y2 = emu_to_px(shape.end_y)
                        col = lcolor or (180, 180, 180)
                        draw.line([(x1, y1), (x2, y2)], fill=col, width=max(1, lwidth))
                    except Exception:
                        pass
                    continue

                # ── Rectangle / auto-shape ────────────────────────────────────
                box = [x, y, x + w, y + h]

                if fill:
                    draw.rectangle(box, fill=fill)

                if lcolor and lwidth > 0:
                    for t in range(min(lwidth, 3)):
                        b2 = [x + t, y + t, x + w - t, y + h - t]
                        draw.rectangle(b2, outline=lcolor)

                # ── Text inside shape ─────────────────────────────────────────
                if not shape.has_text_frame:
                    continue

                tf = shape.text_frame
                full_text = tf.text.strip()
                if not full_text:
                    continue

                # Font size: use first run's size, fallback to 8pt
                font_size_pt = 8
                try:
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                font_size_pt = run.font.size / 12700
                                break
                        else:
                            continue
                        break
                except Exception:
                    pass

                font_px = max(7, int(font_size_pt * SCALE / 72))

                # Text color
                txt_color = (255, 255, 255)
                try:
                    for para in tf.paragraphs:
                        for run in para.runs:
                            c = run.font.color.rgb
                            txt_color = (c[0], c[1], c[2])
                            break
                        else:
                            continue
                        break
                except Exception:
                    # Auto-pick contrast color based on fill
                    if fill:
                        brightness = (fill[0]*299 + fill[1]*587 + fill[2]*114) / 1000
                        txt_color  = (255, 255, 255) if brightness < 140 else (30, 30, 30)
                    else:
                        txt_color = (30, 30, 30)

                # Try to load a font; fall back to default
                try:
                    font = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
                                             font_px)
                except Exception:
                    try:
                        font = ImageFont.truetype('/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
                                                 font_px)
                    except Exception:
                        font = ImageFont.load_default()

                # Truncate text to fit width
                max_chars = max(3, int(w / max(font_px * 0.6, 1)))
                display   = full_text[:max_chars]

                # Center text in shape
                try:
                    bbox   = draw.textbbox((0, 0), display, font=font)
                    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
                except Exception:
                    tw, th = font_px * len(display) * 0.6, font_px

                tx = x + max(0, (w - tw) // 2)
                ty = y + max(0, (h - th) // 2)

                # Clip text to shape bounds
                if tx < x + w and ty < y + h:
                    draw.text((tx, ty), display, fill=txt_color, font=font)

            except Exception:
                continue   # skip any shape that errors

        # ── Encode to PNG bytes ───────────────────────────────────────────────
        buf = _io.BytesIO()
        img.save(buf, format='PNG', optimize=True)
        return buf.getvalue()

    except Exception:
        return None


if go:
    bar = st.progress(0)
    msg = st.empty()
    for pct_val, text in [(20,'Grouping components…'), (50,'Calculating layout…'),
                           (75,'Drawing diagram…'), (90,'Rendering preview…')]:
        msg.markdown(f'<div style="font-family:JetBrains Mono,monospace;font-size:11px;color:#5A6080;">{text}</div>', unsafe_allow_html=True)
        bar.progress(pct_val)

    structure = group_bom(df)
    out = tempfile.mktemp(suffix='.pptx')
    generate_pptx(structure, out, project_title=project_title)

    with open(out, 'rb') as f:
        pptx_bytes = f.read()
    os.unlink(out)

    # Generate preview image
    img_bytes = pptx_to_preview_image(pptx_bytes)

    bar.progress(100)
    bar.empty(); msg.empty()

    st.session_state['pptx_bytes']   = pptx_bytes
    st.session_state['preview_img']  = img_bytes
    st.session_state['preview_title'] = project_title


if st.session_state.get('pptx_bytes'):
    fname     = st.session_state.get('preview_title', project_title).replace(' ','_').replace('/','_') + '_Architecture.pptx'
    img_bytes = st.session_state.get('preview_img')

    # ── PREVIEW PANEL ─────────────────────────────────────────────────────────
    st.markdown("""
    <div class="preview-wrap">
        <div class="preview-topbar">
            <div class="preview-dots">
                <div class="preview-dot preview-dot-r"></div>
                <div class="preview-dot preview-dot-y"></div>
                <div class="preview-dot preview-dot-g"></div>
            </div>
            <div class="preview-label">ARCHITECTURE DIAGRAM · SLIDE 1 OF 1</div>
            <div class="preview-badge">✓ GENERATED</div>
        </div>
        <div class="preview-img-wrap">
    """, unsafe_allow_html=True)

    if img_bytes:
        st.image(img_bytes, use_container_width=True)
    else:
        st.markdown("""
        <div style="padding:3rem;text-align:center;font-family:'JetBrains Mono',monospace;
             font-size:12px;color:#404560;">
            Preview unavailable — download the PPT to view
        </div>
        """, unsafe_allow_html=True)

    st.markdown(f"""
        </div>
        <div class="preview-footer">
            <div class="preview-footer-left">DELTAV SYSTEM ARCHITECTURE · {st.session_state.get('preview_title','').upper()}</div>
            <div class="preview-zoom-hint">DOWNLOAD FOR FULL RESOLUTION</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── DOWNLOAD ROW ──────────────────────────────────────────────────────────
    col_dl, col_info = st.columns([2, 3])
    with col_dl:
        st.download_button(
            '↓  DOWNLOAD ARCHITECTURE.PPTX',
            data=st.session_state['pptx_bytes'],
            file_name=fname,
            mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            use_container_width=True
        )
    with col_info:
        st.markdown(f"""
        <div style="padding:0.7rem 1rem;background:rgba(255,255,255,0.02);
             border:1px solid rgba(255,255,255,0.06);border-radius:8px;
             display:flex;gap:2rem;align-items:center;">
            <div>
                <div style="font-family:'JetBrains Mono',monospace;font-size:9px;
                     letter-spacing:0.12em;color:#404560;text-transform:uppercase;margin-bottom:3px;">Format</div>
                <div style="font-family:'Syne',sans-serif;font-size:13px;
                     font-weight:600;color:#C0C4D8;">.pptx · Microsoft PowerPoint</div>
            </div>
            <div>
                <div style="font-family:'JetBrains Mono',monospace;font-size:9px;
                     letter-spacing:0.12em;color:#404560;text-transform:uppercase;margin-bottom:3px;">Slides</div>
                <div style="font-family:'Syne',sans-serif;font-size:13px;
                     font-weight:600;color:#C0C4D8;">1 · Architecture Only</div>
            </div>
            <div>
                <div style="font-family:'JetBrains Mono',monospace;font-size:9px;
                     letter-spacing:0.12em;color:#404560;text-transform:uppercase;margin-bottom:3px;">Size</div>
                <div style="font-family:'Syne',sans-serif;font-size:13px;
                     font-weight:600;color:#C0C4D8;">{len(st.session_state['pptx_bytes'])//1024} KB</div>
            </div>
        </div>
        """, unsafe_allow_html=True)


# ── FOOTER ────────────────────────────────────────────────────────────────────
st.markdown("""
<div style="margin-top:4rem;padding-top:1.5rem;border-top:1px solid rgba(255,255,255,0.05);
     display:flex;justify-content:space-between;flex-wrap:wrap;gap:0.5rem;">
    <div style="font-family:'JetBrains Mono',monospace;font-size:10px;color:#252840;">DELTAV ARCHITECTURE GENERATOR · EMERSON AUTOMATION SOLUTIONS · PUNE</div>
    <div style="font-family:'JetBrains Mono',monospace;font-size:10px;color:#252840;">RULE-BASED ENGINE · 100% FREE · FULLY OFFLINE</div>
</div>
""", unsafe_allow_html=True)